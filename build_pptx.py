"""Builds Mapping_the_Unknown_Workshop.pptx — a live-teaching deck.

Design standard: on-slide text stays short (scannable while presenting);
full talking-point scripts live in speaker notes for the instructor to read
off while slides are on screen. English leads, Traditional Chinese supports.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0x2E, 0x7D, 0x5B)
GREY = RGBColor(0x6B, 0x6B, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF2, 0xEC)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

TOTAL_MODULES = 8
_slide_no = [0]


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_footer(slide, module_no=None, minutes=None, label=None):
    _slide_no[0] += 1
    box = slide.shapes.add_textbox(Inches(0.5), SH - Inches(0.45), Inches(6), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Mapping the Unknown · 探索未知空間"
    r.font.size = Pt(10)
    r.font.color.rgb = GREY

    if label:
        box2 = slide.shapes.add_textbox(Inches(9.5), SH - Inches(0.45), Inches(3.2), Inches(0.35))
        tf2 = box2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        r2 = p2.add_run()
        r2.text = label
        r2.font.size = Pt(10)
        r2.font.color.rgb = GREY

    numbox = slide.shapes.add_textbox(SW - Inches(0.9), SH - Inches(0.45), Inches(0.6), Inches(0.35))
    ntf = numbox.text_frame
    np = ntf.paragraphs[0]
    np.alignment = PP_ALIGN.RIGHT
    nr = np.add_run()
    nr.text = str(_slide_no[0])
    nr.font.size = Pt(10)
    nr.font.color.rgb = GREY


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_content_slide(kicker, title_en, title_zh, bullets, notes, module_no=None, minutes=None, diagram=None):
    """bullets: list of (en_text, zh_text_or_None, level).
    diagram: optional callable(slide) that draws a static figure in the
    right-hand panel; when set, the bullet column narrows to make room."""
    slide = prs.slides.add_slide(BLANK)
    bg(slide, WHITE)

    if kicker:
        kb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(11.5), Inches(0.4))
        kp = kb.text_frame.paragraphs[0]
        kr = kp.add_run()
        kr.text = kicker
        kr.font.size = Pt(14)
        kr.font.bold = True
        kr.font.color.rgb = ACCENT

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.75), Inches(12.1), Inches(1.15))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title_en
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = NAVY
    if title_zh:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = title_zh
        r2.font.size = Pt(16)
        r2.font.color.rgb = GREY

    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.95), Inches(12.1), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    line.shadow.inherit = False

    body_w = Inches(7.2) if diagram else Inches(12.1)
    body = slide.shapes.add_textbox(Inches(0.6), Inches(2.25), body_w, Inches(4.6))
    btf = body.text_frame
    btf.word_wrap = True
    first = True
    for text, zh, level in bullets:
        p = btf.paragraphs[0] if first else btf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(10 if level == 0 else 4)
        run = p.add_run()
        run.text = ("● " if level == 0 else "– ") + text
        run.font.size = Pt(20) if level == 0 else Pt(16)
        run.font.color.rgb = NAVY if level == 0 else GREY
        if zh:
            p2 = btf.add_paragraph()
            p2.level = level
            p2.space_after = Pt(10 if level == 0 else 4)
            r2 = p2.add_run()
            r2.text = "   " + zh
            r2.font.size = Pt(13) if level == 0 else Pt(11)
            r2.font.italic = True
            r2.font.color.rgb = GREY

    if diagram:
        diagram(slide)

    label = f"Module {module_no:02d}/{TOTAL_MODULES} · {minutes} min" if module_no else None
    add_footer(slide, module_no, minutes, label)
    set_notes(slide, notes)
    return slide


def draw_geometry_diagram(slide):
    """Static point / line / polygon pictograms with an attribute-table icon,
    for the 'From Space to Data' slide."""
    px = Inches(8.15)
    pw = Inches(4.55)
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, Inches(2.25), pw, Inches(4.55))
    panel.fill.solid()
    panel.fill.fore_color.rgb = LIGHT
    panel.line.color.rgb = RGBColor(0xE0, 0xDD, 0xD3)
    panel.line.width = Pt(0.75)
    panel.shadow.inherit = False
    panel.text_frame.margin_left = 0
    panel.text_frame.margin_top = 0

    def row_label(y, text_en, text_zh):
        lb = slide.shapes.add_textbox(px + Inches(2.05), y, pw - Inches(2.2), Inches(0.55))
        ltf = lb.text_frame
        ltf.word_wrap = True
        lp = ltf.paragraphs[0]
        lr = lp.add_run()
        lr.text = text_en
        lr.font.size = Pt(14)
        lr.font.bold = True
        lr.font.color.rgb = NAVY
        lp2 = ltf.add_paragraph()
        lr2 = lp2.add_run()
        lr2.text = text_zh
        lr2.font.size = Pt(11)
        lr2.font.color.rgb = GREY

    icon_cx = px + Inches(1.0)

    # Point — a small filled dot (e.g. a tree)
    y1 = Inches(2.55)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, icon_cx - Inches(0.12), y1 + Inches(0.14), Inches(0.24), Inches(0.24))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()
    dot.shadow.inherit = False
    row_label(y1, "Point — a tree", "點——一棵樹")

    # Line — a street centerline
    y2 = Inches(3.55)
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, icon_cx - Inches(0.55), y2 + Inches(0.28),
                                     icon_cx + Inches(0.55), y2 + Inches(0.05))
    ln.line.color.rgb = ACCENT
    ln.line.width = Pt(3)
    row_label(y2, "Line — a street", "線——一條街道")

    # Polygon — a building footprint
    y3 = Inches(4.55)
    poly = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, icon_cx - Inches(0.45), y3 + Inches(0.08), Inches(0.9), Inches(0.55))
    poly.fill.solid()
    poly.fill.fore_color.rgb = RGBColor(0xBF, 0xDA, 0xCC)
    poly.line.color.rgb = ACCENT
    poly.line.width = Pt(1.5)
    poly.shadow.inherit = False
    row_label(y3, "Polygon — a building", "面——一棟建築")

    # Attributes — a mini table icon
    y4 = Inches(5.65)
    tbl_w, tbl_h = Inches(1.0), Inches(0.62)
    tbl_x, tbl_y = icon_cx - tbl_w / 2, y4 + Inches(0.05)
    tbl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, tbl_x, tbl_y, tbl_w, tbl_h)
    tbl.fill.solid()
    tbl.fill.fore_color.rgb = WHITE
    tbl.line.color.rgb = NAVY
    tbl.line.width = Pt(1.25)
    tbl.shadow.inherit = False
    for gx in (1, 2):
        gl = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, tbl_x + tbl_w * gx / 3, tbl_y,
                                         tbl_x + tbl_w * gx / 3, tbl_y + tbl_h)
        gl.line.color.rgb = NAVY
        gl.line.width = Pt(0.75)
    gl2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, tbl_x, tbl_y + tbl_h / 3,
                                      tbl_x + tbl_w, tbl_y + tbl_h / 3)
    gl2.line.color.rgb = NAVY
    gl2.line.width = Pt(0.75)
    row_label(y4, "Attributes — id, height, use", "屬性——編號、高度、用途")


def add_divider(module_no, title_en, title_zh, minutes, focus_en, focus_zh, notes):
    slide = prs.slides.add_slide(BLANK)
    bg(slide, NAVY)

    numbox = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(3), Inches(1.5))
    ntf = numbox.text_frame
    np = ntf.paragraphs[0]
    nr = np.add_run()
    nr.text = f"{module_no:02d}"
    nr.font.size = Pt(64)
    nr.font.bold = True
    nr.font.color.rgb = ACCENT

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(11.9), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title_en
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = title_zh
    r2.font.size = Pt(22)
    r2.font.color.rgb = RGBColor(0xC9, 0xD3, 0xE0)

    mb = slide.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(11.9), Inches(0.5))
    mp = mb.text_frame.paragraphs[0]
    mr = mp.add_run()
    mr.text = f"{minutes} minutes  ·  {minutes} 分鐘"
    mr.font.size = Pt(16)
    mr.font.color.rgb = ACCENT

    fb = slide.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(11.9), Inches(1.4))
    ftf = fb.text_frame
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fr = fp.add_run()
    fr.text = focus_en
    fr.font.size = Pt(16)
    fr.font.color.rgb = WHITE
    fp2 = ftf.add_paragraph()
    fr2 = fp2.add_run()
    fr2.text = focus_zh
    fr2.font.size = Pt(13)
    fr2.font.color.rgb = RGBColor(0xC9, 0xD3, 0xE0)

    label = f"Module {module_no:02d}/{TOTAL_MODULES} · {minutes} min"
    add_footer(slide, module_no, minutes, label)
    set_notes(slide, notes)
    return slide


# ---------------------------------------------------------------
# 0. Title slide
# ---------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.6))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Mapping the Unknown"
r.font.size = Pt(48)
r.font.bold = True
r.font.color.rgb = WHITE
p2 = tf.add_paragraph()
r2 = p2.add_run()
r2.text = "GIS, Open Data & Spatial Analysis"
r2.font.size = Pt(24)
r2.font.color.rgb = ACCENT
p3 = tf.add_paragraph()
r3 = p3.add_run()
r3.text = "探索未知空間：GIS、開放資料與空間分析"
r3.font.size = Pt(20)
r3.font.color.rgb = RGBColor(0xC9, 0xD3, 0xE0)
p4 = tf.add_paragraph()
r4 = p4.add_run()
r4.text = "A 2-Hour Workshop for Architecture Graduates  ·  建築系畢業生兩小時工作坊"
r4.font.size = Pt(15)
r4.font.italic = True
r4.font.color.rgb = RGBColor(0xC9, 0xD3, 0xE0)
add_footer(s, label="Welcome · 歡迎")
set_notes(s,
    "Welcome learners. Introduce yourself briefly. State the two-hour arc up front: "
    "we move from a research question, through open data and GIS, to a designed research map. "
    "Say the two guiding lines out loud: 'Don't start with a map. Start with a question.' "
    "and 'GIS is not only for making maps — it's a way to think spatially.' "
    "開場：簡短自我介紹，說明兩小時的整體路徑——從研究問題出發，經開放資料與 GIS，最後產出一張研究地圖。"
    "唸出兩句核心理念：「不要從地圖開始，從問題開始」與「GIS 不只是用來畫地圖，而是一種空間思考方式」。")

# ---------------------------------------------------------------
# 1. Mission
# ---------------------------------------------------------------
add_content_slide(
    "Project Mission · 專案宗旨", "Why This Workshop Exists", "為什麼有這個工作坊",
    [
        ("Don't start with a map. Start with a question.", "不要從地圖開始，從問題開始。", 0),
        ("GIS is not only for making maps — it is a way to think spatially.", "GIS 不只是用來畫地圖，而是一種空間思考方式。", 0),
        ("Research Question → Open Data → GIS Exploration → Computational Analysis → Insight → Research Map",
         "研究問題 → 開放資料 → GIS 探索 → 運算分析 → 洞見 → 研究地圖", 0),
    ],
    "Most GIS tutorials start with software. This workshop starts with a research question, "
    "and only reaches for GIS or code when the question demands it. Walk through the chain on screen "
    "left to right — this is the spine of the entire two hours; refer back to it at each module transition. "
    "多數 GIS 教學從軟體操作開始，本工作坊從研究問題出發，只有當問題需要時才引入 GIS 與程式。"
    "這條鏈是整場工作坊的骨幹，之後每個模組轉換時都可以回頭指這張圖。",
)

# ---------------------------------------------------------------
# 2. Agenda
# ---------------------------------------------------------------
add_content_slide(
    "Agenda · 議程", "2-Hour Schedule", "兩小時課程時程",
    [
        ("01  Mapping the Unknown — 10 min", "探索未知空間", 0),
        ("02  From Space to Data — 15 min", "從空間到資料", 0),
        ("03  Finding Open Geospatial Data — 15 min", "尋找開放地理資料", 0),
        ("04  QGIS Basics — 20 min", "QGIS 基礎", 0),
        ("05  Computational GIS with Google Colab — 20 min", "使用 Google Colab 的運算式 GIS", 0),
        ("06  From Data to Spatial Insight — 15 min", "從資料到空間洞察", 0),
        ("07  Research Map Design — 20 min", "研究地圖設計", 0),
        ("08  One Map Challenge — 5 min", "一張地圖挑戰", 0),
    ],
    "Read the agenda aloud once so learners can pace themselves. Mention that modules 01-03 are "
    "conceptual (no software yet), 04-05 are hands-on (QGIS then Colab), 06-07 turn results into "
    "communication, and 08 is the graded deliverable. Point out this slide's structure repeats as a "
    "divider before every module, so learners always know where they are in the two hours. "
    "唸過一次議程讓學員掌握節奏。模組 01-03 是概念性、尚未動手；04-05 動手操作（先 QGIS 後 Colab）；"
    "06-07 把結果轉化為溝通；08 為最終評量產出。之後每個模組前都會有相同格式的分隔頁，讓學員隨時知道進度。",
)

# ---------------------------------------------------------------
# 3. Target audience
# ---------------------------------------------------------------
add_content_slide(
    "Before We Start · 開始之前", "Who This Is For", "適合對象",
    [
        ("Architecture graduates & early-stage spatial researchers", "建築系畢業生與初階空間研究者", 0),
        ("Assumes: basic architectural / spatial literacy", "假設具備：基本建築／空間素養", 0),
        ("Assumes: NO advanced GIS or Python knowledge", "假設不需具備：進階 GIS 或 Python 知識", 0),
        ("Taught through architecture, urban space, landscape, public space, environment",
         "透過建築、都市空間、地景、公共空間、環境案例教學", 0),
    ],
    "Reassure the room: nobody needs prior GIS or coding experience. Every technical concept in this "
    "deck is anchored to something architects already reason about intuitively — site reading, figure-ground, "
    "diagramming. This lowers anxiety before QGIS and Colab show up in modules 04-05. "
    "安撫學員：不需要先備 GIS 或程式經驗。每個技術概念都會對應到建築人早已熟悉的直覺——基地閱讀、圖底關係、圖表化。"
    "這能降低模組 04-05 出現 QGIS 與 Colab 時的焦慮感。",
)

# =================================================================
# MODULE 01 — Mapping the Unknown
# =================================================================
add_divider(1, "Mapping the Unknown", "探索未知空間", 10,
            "What is GIS? What is spatial thinking? GIS as a research method.",
            "什麼是 GIS？什麼是空間思考？GIS 作為研究方法。",
            "Open with the question 'what do you think GIS is?' and take two or three answers — most will "
            "say 'mapping software.' Use that to set up the central reframe of the whole workshop. "
            "以「你們覺得 GIS 是什麼？」開場，收兩三個回答——多數人會說「製圖軟體」。用這個回答帶出整場工作坊的核心翻轉。")

add_content_slide(
    "01 · Mapping the Unknown", "What Is GIS?", "什麼是 GIS？",
    [
        ("A system for capturing, storing, analyzing, and communicating data that has a location",
         "一套擷取、儲存、分析並傳達具有位置的資料的系統", 0),
        ("The software is not the point — what location-aware data lets you ask and answer is",
         "軟體本身不是重點，重點是位置資料能讓你提出並回答什麼問題", 0),
        ("GIS ≠ map making", "GIS 不等於製圖", 0),
        ("Mapping is often the output; GIS is the reasoning that gets you there",
         "製圖經常只是產出；GIS 是抵達產出之前的推理過程", 0),
    ],
    "Emphasize the reframe hard: GIS ≠ map making. Making a map is an output; GIS the discipline is "
    "about the reasoning beforehand — what varies across space, why, and what it means. This single "
    "sentence is the thesis statement for the entire two hours. "
    "強力強調這個翻轉：GIS 不等於製圖。製圖只是產出，GIS 這門學科關注的是產出之前的推理——"
    "什麼東西在空間中變化、為什麼變化、代表什麼意義。這句話是整場工作坊的核心論點。",
    module_no=1, minutes=10,
)

add_content_slide(
    "01 · Mapping the Unknown", "What Is Spatial Thinking?", "什麼是空間思考？",
    [
        ("Location — where something is", "位置——某事物在哪裡", 0),
        ("Distribution — how it's spread out", "分布——如何散布", 0),
        ("Pattern — clustered, even, or random?", "樣式——聚集、均勻或隨機？", 0),
        ("Relationship — how it relates to nearby things", "關係——與鄰近事物的關聯", 0),
        ("Architects already do this reading a site: sun, foot traffic, noise",
         "建築人讀基地時早已這樣做：陽光、人流、噪音", 0),
    ],
    "Land this on the architecture analogy: reading a site for sun path, pedestrian desire lines, and "
    "noise sources is already spatial thinking. GIS just formalizes that intuition into data you can "
    "measure, compare, and communicate — it's not a new skill, it's a new vocabulary for an old skill. "
    "落在建築類比上：解讀基地的陽光、人行動線、噪音來源，本身就是空間思考。GIS 只是把這種直覺"
    "形式化為可測量、可比較、可傳達的資料——不是新技能，是舊技能的新詞彙。",
    module_no=1, minutes=10,
)

# =================================================================
# MODULE 02 — From Space to Data
# =================================================================
add_divider(2, "From Space to Data", "從空間到資料", 15,
            "Vector vs raster, point/line/polygon, attributes, coordinates, CRS.",
            "向量與網格、點線面、屬性資料、座標、CRS。",
            "Transition line: 'Spatial reality is continuous; data is not.' This module is about what gets "
            "lost and gained when we encode architectural space as data. Keep the pace brisk — it's mostly "
            "vocabulary the room needs before QGIS. "
            "過渡句：「真實空間是連續的，資料卻不是。」本模組談論把建築空間編碼為資料時，失去與獲得了什麼。"
            "步調保持明快——這主要是進入 QGIS 前需要的詞彙。")

add_content_slide(
    "02 · From Space to Data", "Vector vs. Raster", "向量與網格",
    [
        ("Vector — discrete geometries: points, lines, polygons", "向量——離散幾何：點、線、面", 0),
        ("Example: a building footprint drawn as a polygon", "範例：以多邊形繪製的建築足跡", 1),
        ("Raster — a grid of cells/pixels, each with a value", "網格——由格網／像元組成，每格皆有數值", 0),
        ("Example: a satellite image, a DEM (elevation raster)", "範例：衛星影像、DEM（高程網格）", 1),
        ("Vector suits discrete objects; raster suits continuous phenomena",
         "向量適合離散物件；網格適合連續現象", 0),
    ],
    "Draw the line clearly: vector for discrete objects (buildings, roads, parcels), raster for continuous "
    "phenomena (temperature, elevation, land cover). Ask the room: is a park boundary vector or raster? "
    "(vector — it's a discrete polygon). Is a heat map of the city raster or vector? (raster). "
    "清楚劃分：向量適合離散物件（建築、道路、地籍），網格適合連續現象（溫度、高程、地表覆蓋）。"
    "問學員：公園邊界是向量還是網格？（向量——離散多邊形）。城市熱力圖是網格還是向量？（網格）。",
    module_no=2, minutes=15,
)

add_content_slide(
    "02 · From Space to Data", "Point, Line, Polygon & Attributes", "點、線、面與屬性資料",
    [
        ("Point — a tree, a lamppost, a building entrance", "點——一棵樹、一根路燈、一個建築入口", 0),
        ("Line — a street centerline, a walking path, a river", "線——街道中心線、步行路徑、河流", 0),
        ("Polygon — a building footprint, a park boundary, a parcel", "面——建築足跡、公園邊界、地籍範圍", 0),
        ("Attributes — the non-spatial facts attached to a geometry", "屬性——附加於幾何物件上的非空間事實資料", 0),
        ("Geometry answers \"where\"; attributes answer \"what\"", "幾何回答「在哪裡」；屬性回答「是什麼」", 0),
    ],
    "Show the attribute-table example verbally: a building footprint polygon carrying id, height_m, "
    "year_built, use. Land the key line: geometry answers 'where,' attributes answer 'what.' Both are "
    "needed to answer a research question — geometry alone or attributes alone are not enough. Point at "
    "the diagram panel on the right while naming each geometry type. "
    "口頭示範屬性表範例：建築足跡多邊形帶有 id、height_m、year_built、use 等欄位。"
    "點出關鍵句：幾何回答「在哪裡」，屬性回答「是什麼」。兩者缺一不可才能回答研究問題。"
    "唸到每種幾何類型時，可指向右側簡圖對應說明。",
    module_no=2, minutes=15,
    diagram=draw_geometry_diagram,
)

add_content_slide(
    "02 · From Space to Data", "Coordinates & CRS", "座標與 CRS",
    [
        ("Coordinates locate a point on Earth's surface (e.g. lat/lon)", "座標定位地表上的一個點（例如經緯度）", 0),
        ("CRS (Coordinate Reference System) defines how those numbers map to real space",
         "CRS（座標參照系統）定義這些數字如何對應到真實空間", 0),
        ("Two layers in different CRS will not align — even if the data is correct",
         "兩圖層若 CRS 不同，即使資料本身正確也無法對齊", 0),
        ("Rule of thumb: always check CRS before comparing or combining layers",
         "經驗法則：疊圖或比對前，一定先檢查 CRS", 0),
    ],
    "This is the single most common beginner error, and it will resurface in Module 04 (QGIS) as "
    "step 3. Plant the seed now: mismatched CRS is why two correct datasets can look wrong together. "
    "You don't need to explain projections in depth — just that CRS is a handshake that must match. "
    "這是初學者最常見的錯誤，會在模組 04（QGIS）第 3 步再次出現。現在先埋下伏筆：CRS 不一致會讓"
    "兩份正確的資料疊在一起卻對不上。不需深入講解投影法，只需說明 CRS 是必須一致的「握手協議」。",
    module_no=2, minutes=15,
)

# =================================================================
# MODULE 03 — Finding Open Geospatial Data
# =================================================================
add_divider(3, "Finding Open Geospatial Data", "尋找開放地理資料", 15,
            "OpenStreetMap, government open data, APIs — worked backwards from your question.",
            "OpenStreetMap、政府開放資料、API——從問題反推資料。",
            "Key behavioral rule for this module: never let students browse a data portal aimlessly. "
            "Everything here works backwards from a question, through required variables, to a dataset. "
            "本模組的關鍵行為準則：絕不讓學員漫無目的瀏覽資料平台。一切都從問題出發，"
            "推導所需變數，再找到資料集。")

add_content_slide(
    "03 · Finding Open Geospatial Data", "The Logic: Work Backwards From the Question", "邏輯：從問題反推資料",
    [
        ("Research Question", "研究問題", 0),
        ("→ Required Variables", "→ 所需變數", 1),
        ("→ Potential Dataset", "→ 可能的資料集", 1),
        ("→ Data Source", "→ 資料來源", 1),
        ("→ Spatial Dataset", "→ 空間資料集", 1),
        ("Example: \"Where are the greener public spaces?\" → tree locations + public space boundaries",
         "範例：「哪裡是較綠意盎然的公共空間？」→ 樹木位置＋公共空間邊界", 0),
    ],
    "Walk the example live: question = 'where are the greener public spaces?'; required variables = "
    "tree locations and public space boundaries; potential dataset = street tree inventory and park/plaza "
    "polygons; source = city open data portal or OSM; final spatial dataset = trees.geojson + "
    "public_spaces.geojson. This exact example returns in Module 08's One Map Challenge. "
    "現場走過範例：問題＝「哪裡是較綠意盎然的公共空間？」；所需變數＝樹木位置與公共空間邊界；"
    "可能資料集＝行道樹清冊與公園／廣場面資料；來源＝城市開放資料平台或 OSM；"
    "最終空間資料集＝trees.geojson＋public_spaces.geojson。這個範例會在模組 08 一張地圖挑戰再次出現。",
    module_no=3, minutes=15,
)

add_content_slide(
    "03 · Finding Open Geospatial Data", "Where to Look", "去哪裡找",
    [
        ("OpenStreetMap (OSM) — free, editable, global; buildings, roads, land use, POIs",
         "OpenStreetMap——免費、可編輯、全球性；建築、道路、土地利用、興趣點", 0),
        ("Government open data portals — city/national datasets", "政府開放資料平台——市級／國家級資料集", 0),
        ("APIs — programmatic access to live or large datasets", "API——以程式化方式存取即時或大型資料集", 0),
        ("Environmental, transport, and weather data sources", "環境、交通與氣象資料來源", 0),
    ],
    "Mention osmnx as the Python bridge to OSM (used in Module 05). Keep this slide as an index, not a "
    "deep dive — the full reference list lives in resources/open-data/ for learners to revisit later. "
    "Note in passing that real Taiwan data sources (data.taipei, Open-Meteo, CWA opendata) power the "
    "optional case studies referenced after Module 05. "
    "提及 osmnx 作為存取 OSM 的 Python 工具（模組 05 會用到）。這頁只是索引，不深入——"
    "完整參考清單在 resources/open-data/ 供學員之後查閱。順帶提及真實台灣資料來源"
    "（data.taipei、Open-Meteo、中央氣象署開放資料）會用於模組 05 之後的延伸案例研究。",
    module_no=3, minutes=15,
)

# =================================================================
# MODULE 04 — QGIS Basics
# =================================================================
add_divider(4, "QGIS Basics", "QGIS 基礎", 20,
            "Load, inspect, filter, buffer, spatial join, intersect, symbolize.",
            "載入、檢視、篩選、緩衝區、空間 join、交集、符號化。",
            "Hands-on begins here. State the module's guiding rule up front: every operation follows "
            "Tool → Spatial Meaning → Research Question. Never let a click stand without its 'why.' "
            "動手操作從這裡開始。先講清楚本模組的準則：每個操作都遵循 工具 → 空間意義 → 研究問題。"
            "絕不讓一次點擊沒有「為什麼」。")

add_content_slide(
    "04 · QGIS Basics", "Load, Inspect, Check CRS", "載入、檢視、檢查 CRS",
    [
        ("1. Load data — drag a GeoJSON or CSV with lat/lon into QGIS", "1. 載入資料——將 GeoJSON 或含經緯度的 CSV 拖入 QGIS", 0),
        ("2. Inspect attributes — does this dataset contain the variable my question needs?",
         "2. 檢視屬性——這份資料是否包含問題所需的變數？", 0),
        ("3. Check CRS — confirm all layers share a CRS before anything else",
         "3. 檢查 CRS——先確認所有圖層 CRS 一致，再做其他事", 0),
        ("Misaligned layers are the #1 beginner error", "圖層對不齊是初學者最常見的錯誤", 0),
    ],
    "This is where the CRS warning from Module 02 pays off — call back to it explicitly: 'remember the "
    "handshake rule from ten minutes ago? Here's where it bites you if you skip it.' Loading data is the "
    "moment a research question first meets actual geometry. "
    "這裡正是模組 02 CRS 提醒的回收點——明確呼應：「還記得十分鐘前的握手規則嗎？"
    "如果跳過它，這裡就會出問題。」載入資料是研究問題首次與實際幾何相遇的時刻。",
    module_no=4, minutes=20,
)

add_content_slide(
    "04 · QGIS Basics", "Filter, Buffer, Spatial Join, Intersection", "篩選、緩衝區、空間 Join、交集",
    [
        ("4. Filter — isolate features relevant to the question", "4. 篩選——篩出與問題相關的要素", 0),
        ("5. Buffer — turn a distance question (\"within 10 m\") into a shape",
         "5. 緩衝區——將距離問題（「10 公尺內」）轉換為形狀", 0),
        ("6. Spatial Join — attach attributes from one layer to another by location",
         "6. 空間 Join——依位置將一圖層的屬性附加到另一圖層", 0),
        ("7. Intersection — find the exact overlap between two layers", "7. 交集——找出兩圖層間精確的重疊範圍", 0),
    ],
    "Read the canonical example aloud for Buffer: 'We create a 10 m buffer because we want to ask which "
    "buildings are located within 10 m of vegetation.' Never present a tool without its research-question "
    "framing — that is the pedagogical rule for this entire repository. "
    "唸出緩衝區的經典範例：「我們建立 10 公尺緩衝區，是因為想知道哪些建築位於植栽 10 公尺範圍內。」"
    "絕不在沒有研究問題框架的情況下介紹工具——這是整個教材的教學準則。",
    module_no=4, minutes=20,
)

add_content_slide(
    "04 · QGIS Basics", "Symbology & The Real Point", "符號化與核心重點",
    [
        ("8. Basic symbology — color or size features by attribute value",
         "8. 基本符號化——依屬性值為要素上色或調整大小", 0),
        ("Memorize the question each tool answers, not the menu path",
         "記住每個工具回答的問題，而非選單路徑", 0),
        ("Filter → which features match a condition?", "篩選 → 哪些要素符合條件？", 1),
        ("Buffer → what's within X distance?", "緩衝區 → 什麼在 X 距離內？", 1),
        ("Spatial Join → what does this location inherit from its context?", "空間 Join → 此位置從脈絡繼承了什麼？", 1),
        ("Intersection → where do two things overlap?", "交集 → 兩者哪裡重疊？", 1),
    ],
    "Close the module with this table as a memory anchor — it's more durable than menu screenshots. "
    "Bridge to Module 05: 'Now do these exact same four operations in code, at scale, reproducibly.' "
    "以這張對照表為記憶錨點作結——比選單截圖更持久。銜接模組 05："
    "「接下來用程式碼、規模化、可重現的方式，執行完全相同的四個操作。」",
    module_no=4, minutes=20,
)

# =================================================================
# MODULE 05 — Computational GIS with Google Colab
# =================================================================
add_divider(5, "Computational GIS with Google Colab", "使用 Google Colab 的運算式 GIS", 20,
            "Python, GeoPandas, Shapely, Matplotlib — the same operations, reproducibly.",
            "Python、GeoPandas、Shapely、Matplotlib——相同操作，但可重現。",
            "Set expectations: this is not a programming course. The code is beginner-friendly and "
            "runs entirely in the browser via Colab — nobody installs anything locally. "
            "先設定期望：這不是一門程式設計課。程式碼對初學者友善，完全透過瀏覽器上的 Colab 執行——"
            "沒有人需要在本機安裝任何東西。")

add_content_slide(
    "05 · Computational GIS", "Why Code, If QGIS Already Works?", "QGIS 已經能用，為何還要寫程式？",
    [
        ("QGIS = See the Space — visual exploration, spatial intuition", "QGIS ＝看見空間——視覺探索、建立空間直覺", 0),
        ("Colab / Python = Compute the Space — reproducible, scalable operations",
         "Colab／Python ＝運算空間——可重現、可擴充的空間運算", 0),
        ("Same four operations (filter, buffer, join, intersect) — now scriptable and repeatable",
         "同樣四個操作（篩選、緩衝區、join、交集）——現在可寫成腳本並重複執行", 0),
        ("Tools: Python, GeoPandas, Pandas, Shapely, Matplotlib", "工具：Python、GeoPandas、Pandas、Shapely、Matplotlib", 0),
    ],
    "Frame Colab as QGIS's complement, not its replacement. QGIS builds intuition by seeing; Colab scales "
    "that same reasoning into something reproducible — rerun the whole analysis on new data in seconds "
    "instead of re-clicking through a GUI. "
    "把 Colab 定位為 QGIS 的互補，而非取代。QGIS 透過「看見」建立直覺；Colab 把同樣的推理規模化為"
    "可重現的流程——在新資料上重跑整個分析只需幾秒，不必在圖形介面裡重新點擊一次。",
    module_no=5, minutes=20,
)

add_content_slide(
    "05 · Computational GIS", "The Workflow", "工作流程",
    [
        ("Load → Inspect → Clean → Spatial Operation → Calculate → Visualize",
         "載入 → 檢視 → 清理 → 空間運算 → 計算 → 視覺化", 0),
        ("Every notebook runs entirely in Google Colab — no local install", "每份筆記本完全在 Google Colab 執行——無需本機安裝", 0),
        ("First cell installs dependencies automatically (geopandas, shapely, matplotlib)",
         "第一個儲存格會自動安裝所需套件", 0),
        ("Open in Colab badge on every notebook", "每份筆記本皆附「Open in Colab」徽章", 0),
    ],
    "Demo live if possible: open notebooks/01_open_data.ipynb via its Colab badge and run the first two "
    "cells so the room sees the install-and-load pattern once. If short on time, screenshot instead — "
    "the goal is just to demystify what 'running a notebook' looks like. "
    "若時間允許，現場示範：透過 Colab 徽章開啟 notebooks/01_open_data.ipynb，執行前兩個儲存格，"
    "讓學員看一次安裝與載入的流程。時間不夠可改用截圖——目的只是讓「執行筆記本」不再神秘。",
    module_no=5, minutes=20,
)

add_content_slide(
    "05 · Computational GIS", "Want Real Data? Case Studies", "想用真實資料？案例研究",
    [
        ("Core notebooks use small hand-made sample data on purpose — no network stalls in class",
         "核心筆記本刻意使用小型手造範例資料——課堂上不會卡在網路問題", 0),
        ("Case Study 01 — Green Coverage Grid: real Taipei street-tree data, no registration",
         "案例一——綠覆率網格化：真實台北市行道樹資料，無需註冊", 0),
        ("Case Study 02 — Urban Heat Interpolation: Open-Meteo API, spatial interpolation",
         "案例二——都市氣溫熱力圖：Open-Meteo API，空間內插", 0),
        ("Recommended as the real-data option for this module's demo and the One Map Challenge",
         "建議作為本模組示範與一張地圖挑戰的真實資料選項", 0),
    ],
    "These are optional, self-paced extensions living in case-studies/ — point learners there for "
    "homework or their final project, but do not attempt to run them live unless time is generous, "
    "since Case Study 02's advanced path needs a free API key. Case Study 01 needs no registration at "
    "all and is the safest live-demo choice. "
    "這些是選用的自學延伸內容，位於 case-studies/ ——可指引學員作為課後作業或期末專案的起點，"
    "但除非時間充裕，否則不建議現場執行，因為案例二的進階路徑需要免費 API 金鑰。"
    "案例一完全無需註冊，是最安全的現場示範選項。",
    module_no=5, minutes=20,
)

# =================================================================
# MODULE 06 — From Data to Spatial Insight
# =================================================================
add_divider(6, "From Data to Spatial Insight", "從資料到空間洞察", 15,
            "A computed result is not an insight. What? Where? Why? So what?",
            "運算結果不是洞見。是什麼？在哪裡？為什麼？所以呢？",
            "Pace shift: close the laptops here if possible. This module is about interpretation, not "
            "tools — it works better as discussion than demonstration. "
            "節奏轉換：這裡可以請學員先闔上筆電。本模組談的是詮釋而非工具——用討論比示範更有效。")

add_content_slide(
    "06 · Spatial Insight", "The Bridge: Result → Insight", "橋樑：結果 → 洞見",
    [
        ("Data → Pattern → Relationship → Interpretation → Insight", "資料 → 樣式 → 關係 → 詮釋 → 洞見", 0),
        ("\"I ran an analysis\" is not the same as \"I learned something\"",
         "「我完成了一次分析」不等於「我發現了某件事」", 0),
        ("Push every result through four questions, in order", "把每個結果依序推過四個問題", 0),
    ],
    "State plainly: a computed result is not an insight. This module is the bridge between running an "
    "analysis and actually learning something from it. Introduce the four-question funnel before showing "
    "the worked example on the next slide. "
    "直接點明：運算結果本身不是洞見。本模組是「完成分析」與「真正學到東西」之間的橋樑。"
    "先介紹四個問題的漏斗架構，再進入下一張的範例。",
    module_no=6, minutes=15,
)

add_content_slide(
    "06 · Spatial Insight", "What? Where? Why? So What?", "是什麼？在哪裡？為什麼？所以呢？",
    [
        ("What? — describe the pattern, don't interpret yet", "是什麼？——先描述樣式，還不要詮釋", 0),
        ("\"60% of public spaces have no tree within 10m; 40% do.\"", "「60% 的公共空間 10 公尺內無樹木；40% 有。」", 1),
        ("Where? — where does it concentrate, cluster, or vary?", "在哪裡？——集中、聚集或變化於何處？", 0),
        ("Why? — a grounded hypothesis, not yet proven", "為什麼？——根植於認識的假設，尚未證實", 0),
        ("So what? — why it matters for the research question, design, policy",
         "所以呢？——為何這對研究問題、設計、政策有意義", 0),
    ],
    "Run this exact example live: 'the tree-poor spaces cluster in the older, denser district' (Where), "
    "'that district was built before street-tree planting requirements existed' (Why — flag as hypothesis, "
    "not fact). Close with: the map is evidence for an argument, not the conclusion itself. "
    "現場走過這個範例：「缺乏樹木的空間集中於較舊、密度較高的區域」（在哪裡），"
    "「該區域是在行道樹種植規範出現之前建成的」（為什麼——標明這是假設，非事實）。"
    "以此作結：地圖是論證的證據，而不是結論本身。",
    module_no=6, minutes=15,
)

# =================================================================
# MODULE 07 — Research Map Design
# =================================================================
add_divider(7, "Research Map Design", "研究地圖設計", 20,
            "Hierarchy, figure-ground, color, typography, annotation, composition.",
            "層級、圖底關係、色彩、字體排印、註記、構圖。",
            "This is the longest module and the one architecture graduates usually enjoy most — it's "
            "their home turf. Move briskly through terms they already know (hierarchy, figure-ground, "
            "composition) and slow down on color encoding, which is usually new. "
            "這是最長的模組，也通常是建築系畢業生最享受的部分——這是他們的主場。"
            "熟悉的詞彙（層級、圖底、構圖）可以講快一點，色彩編碼通常是新概念，要放慢。")

add_content_slide(
    "07 · Research Map Design", "Visual Hierarchy & Figure-Ground", "視覺層級與圖底關係",
    [
        ("What should the viewer see first? Decide this deliberately",
         "觀者應該先看到什麼？需刻意決定", 0),
        ("Usually the main spatial pattern — not the title or legend", "通常是主要空間樣式——而非標題或圖例", 0),
        ("Figure-ground — borrowed from architectural drawing", "圖底關係——直接借用自建築製圖", 0),
        ("What is figure (the subject of analysis)? What is ground (context)?",
         "什麼是圖（分析主體）？什麼是底（脈絡背景）？", 0),
        ("A research map should draw this line as deliberately as a Nolli map",
         "研究地圖應如 Nolli 地圖般刻意做出這個區分", 0),
    ],
    "The Nolli map reference lands well with architecture graduates — most will recognize it. Use it as "
    "the anchor image: figure-ground was never a foreign GIS concept, it's their own discipline's tool "
    "pointed at a new kind of drawing. "
    "Nolli 地圖的類比對建築系畢業生特別有效——多數人會認得。用它作為錨點意象："
    "圖底關係從來不是外來的 GIS 概念，而是他們自己學科的工具，指向一種新的圖面。",
    module_no=7, minutes=20,
)

add_content_slide(
    "07 · Research Map Design", "Color Must Encode Information", "色彩必須傳達資訊",
    [
        ("Color encodes — it never just decorates", "色彩用來編碼——絕不只是裝飾", 0),
        ("Sequential — ordered, low to high (e.g. tree density: light → dark green)",
         "連續型——有序、由低到高（例如樹木密度：淺綠到深綠）", 0),
        ("Diverging — a meaningful midpoint (e.g. above/below average temperature)",
         "發散型——具有意義的中點（例如高於／低於平均溫度）", 0),
        ("Categorical — unordered categories (e.g. land use type)", "類別型——無序類別（例如土地利用類型）", 0),
    ],
    "Quiz the room quickly: 'land use type — sequential, diverging, or categorical?' (categorical — no "
    "inherent order). 'Temperature deviation from the city average?' (diverging — zero is meaningful). "
    "Picking the wrong family is the most common map-design mistake in the One Map Challenge submissions. "
    "快速考學員：「土地利用類型——連續、發散還是類別型？」（類別型——沒有固有順序）。"
    "「與城市平均溫差？」（發散型——零這個中點有意義）。選錯色彩類型是一張地圖挑戰"
    "作業中最常見的錯誤。",
    module_no=7, minutes=20,
)

add_content_slide(
    "07 · Research Map Design", "Typography, Annotation & Composition", "字體排印、註記與構圖",
    [
        ("Typography order: Title → Subtitle → Legend → Annotation → Data source → Scale → North arrow",
         "字體排印順序：標題 → 副標 → 圖例 → 註記 → 資料來源 → 比例尺 → 指北針", 0),
        ("Annotation — point directly at the finding, one sentence beats a paragraph",
         "註記——直接指出發現，一句話勝過一整段文字", 0),
        ("Composition: hierarchy, whitespace, alignment, contrast, legibility",
         "構圖：層級、留白、對齊、對比、易讀性", 0),
        ("Recommended layout: Title / Main Map / Legend + Key Insight / Source-Scale-North-Author",
         "建議版面：標題／主地圖／圖例＋關鍵洞見／來源－比例尺－指北針－作者", 0),
    ],
    "Show the recommended layout diagram from the README if you have it on hand (title block, main map, "
    "legend paired with a key-insight callout, footer strip with source/scale/north/author). A4 landscape "
    "or 16:9. This exact layout is what the One Map Challenge rubric checks against next. "
    "若手邊有 README 的建議版面圖可以展示（標題區、主地圖、圖例與關鍵洞見並列、含來源／比例尺／"
    "指北針／作者的頁尾）。A4 橫向或 16:9。這正是接下來一張地圖挑戰評分標準比對的版面。",
    module_no=7, minutes=20,
)

add_content_slide(
    "07 · Research Map Design", "Before You Submit: The Checklist", "提交前：檢查清單",
    [
        ("Does the map answer a question?", "地圖是否回答問題？", 0),
        ("Is the main spatial pattern obvious? Is there a clear visual hierarchy?",
         "主要空間樣式是否明顯？視覺層級是否清楚？", 0),
        ("Is color encoding meaningful? Is the legend understandable?",
         "色彩編碼是否有意義？圖例是否易懂？", 0),
        ("Is unnecessary information removed? Is the study area clear?",
         "是否移除多餘資訊？研究範圍是否清楚？", 0),
        ("Is the data source documented? Is the key insight visible?",
         "資料來源是否標註？關鍵洞見是否清晰可見？", 0),
    ],
    "This checklist (resources/map-design/research-map-checklist.md) is the actual grading instrument "
    "for the One Map Challenge — tell learners explicitly that this slide is their rubric. Suggest they "
    "run their own map through it before calling it done. "
    "這份檢查清單（resources/map-design/research-map-checklist.md）就是一張地圖挑戰實際的評分依據——"
    "明確告訴學員這張投影片就是他們的評分標準。建議他們完成地圖後親自跑過一次檢查清單。",
    module_no=7, minutes=20,
)

# =================================================================
# MODULE 08 — One Map Challenge
# =================================================================
add_divider(8, "One Map Challenge", "一張地圖挑戰", 5,
            "1 Research Question + 1 Spatial Analysis + 1 Research Map + 1 Spatial Insight.",
            "1 個研究問題 ＋ 1 項空間分析 ＋ 1 張研究地圖 ＋ 1 個空間洞察。",
            "This is the deliverable. Keep this section fast and concrete — learners should leave "
            "knowing exactly what to produce and where to find help. "
            "這是最終產出。這部分要快而具體——學員離開時應該清楚知道要產出什麼、去哪裡找協助。")

add_content_slide(
    "08 · One Map Challenge", "The Final Deliverable", "最終產出",
    [
        ("1 Research Question", "1 個研究問題", 0),
        ("+ 1 Spatial Analysis", "＋ 1 項空間分析", 0),
        ("+ 1 Research Map", "＋ 1 張研究地圖", 0),
        ("+ 1 Spatial Insight", "＋ 1 個空間洞察", 0),
        ("Example: \"Where are the greener public spaces?\" → tree density within public spaces",
         "範例：「哪裡是較綠意盎然的公共空間？」→ 公共空間內的樹木密度", 0),
    ],
    "Close the loop: this is the exact example first introduced in Module 03. Learners can either reuse "
    "it directly, or apply the same four-part structure to their own question. Point again to "
    "case-studies/01-green-coverage-grid/ as the recommended real-data starting point, and to the "
    "checklist from the previous module as the grading rubric. "
    "收尾呼應：這正是模組 03 一開始介紹的範例。學員可以直接沿用，也可以把同樣的四段式結構"
    "套用到自己的問題上。再次指向 case-studies/01-green-coverage-grid/ 作為建議的真實資料起點，"
    "以及上一模組的檢查清單作為評分依據。",
    module_no=8, minutes=5,
)

# ---------------------------------------------------------------
# Closing: QGIS + Colab + Map relationship recap
# ---------------------------------------------------------------
add_content_slide(
    "Recap · 回顧", "Two Tools, One Workflow", "兩個工具，一條工作流",
    [
        ("QGIS = See the Space", "QGIS ＝看見空間", 0),
        ("Colab / Python = Compute the Space", "Colab／Python ＝運算空間", 0),
        ("Research Map = Explain the Space", "研究地圖＝解釋空間", 0),
        ("Not competitors — see, then compute, then explain", "並非互斥——先看見、再運算、最後解釋", 0),
    ],
    "Use this as the closing recap before the final slide — it compresses the whole two hours into one "
    "diagram. Ask the room to repeat the three verbs back: see, compute, explain. "
    "作為最後投影片前的收尾回顧——這張圖把整整兩小時濃縮成一張圖。請學員複誦這三個動詞："
    "看見、運算、解釋。",
)

# ---------------------------------------------------------------
# Closing slide
# ---------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(2.8))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Map the Unknown → Analyze the Space → Discover the Insight"
r.font.size = Pt(28)
r.font.bold = True
r.font.color.rgb = WHITE
p2 = tf.add_paragraph()
r2 = p2.add_run()
r2.text = "探索未知 → 分析空間 → 發現洞察"
r2.font.size = Pt(20)
r2.font.color.rgb = ACCENT
p3 = tf.add_paragraph()
r3 = p3.add_run()
r3.text = ("\"I can start with a spatial question, find data, analyze it, and communicate what I "
           "discovered through a well-designed research map.\"")
r3.font.size = Pt(15)
r3.font.italic = True
r3.font.color.rgb = RGBColor(0xC9, 0xD3, 0xE0)
p4 = tf.add_paragraph()
p4.space_before = Pt(20)
r4 = p4.add_run()
r4.text = "Questions? Discussion · 問題與討論"
r4.font.size = Pt(16)
r4.font.color.rgb = WHITE
add_footer(s, label="Thank you · 謝謝")
set_notes(s,
    "Read the final quote aloud as the target outcome for every learner in the room. Open the floor for "
    "questions. If time remains, walk to a laptop and let one volunteer try loading their own question "
    "into the Module 03 logic chain live. "
    "把最後這句話唸出來，作為每位學員應達成的目標。開放提問。若還有時間，"
    "可以走到電腦前，讓一位自願學員現場把自己的問題套進模組 03 的邏輯鏈試試看。")

out_path = "Mapping_the_Unknown_Workshop.pptx"
prs.save(out_path)
print("Saved:", out_path, "| slides:", len(prs.slides._sldIdLst))
